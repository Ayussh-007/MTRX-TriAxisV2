"""
MTRX-TriAxis | Lesson Slide Generator  (v3 — no-overlap layout)
Uses RAG + Ollama LLM to build rich, classroom-ready slide content,
then renders a branded .pptx with proper paragraph spacing so text
never overlaps even when sentences wrap across multiple lines.

Layout principle
────────────────
  • All bullet lists live in ONE textframe with add_paragraph().
    python-pptx handles line-flow; no fixed-Y-per-bullet.
  • Slide canvas: 13.33" × 7.5"  (widescreen standard)
  • Title band : 0  – 1.05"
  • Body area  : 1.10" – 6.60"  (5.5" tall, split into cols)
  • Footer     : 6.65" – 7.5"
"""

import io
import json
import re
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from backend.llm_utils import get_llm
from backend.rag_pipeline import retrieve_context, vectorstore_exists, load_vectorstore


# ── Palette ───────────────────────────────────────────────────────────────────
CLR_BG_DARK   = RGBColor(0x0E, 0x11, 0x17)
CLR_BG_CARD   = RGBColor(0x1A, 0x1D, 0x29)
CLR_BG_PANEL  = RGBColor(0x10, 0x13, 0x1E)
CLR_ACCENT    = RGBColor(0x7C, 0x6F, 0xFF)
CLR_ACCENT2   = RGBColor(0xA7, 0x8B, 0xFA)
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLR_LGREY     = RGBColor(0xC4, 0xB5, 0xFD)
CLR_GREY      = RGBColor(0x6B, 0x72, 0x80)
CLR_SILVER    = RGBColor(0xD1, 0xD5, 0xDB)
CLR_ORANGE    = RGBColor(0xFB, 0x92, 0x3C)
CLR_GREEN     = RGBColor(0x22, 0xC5, 0x5E)
CLR_BLUE      = RGBColor(0x60, 0xA5, 0xFA)

# Slide canvas (widescreen 16:9)
W = 13.33   # inches
H = 7.50    # inches


# ═══════════════════════════════════════════════════════════════════════════════
# LLM PROMPT + CONTENT GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def _build_prompt(topic, context, num_content, style, difficulty,
                  examples, key_terms, quiz_slide):

    extras, example_note = [], ""
    if key_terms:
        extras.append("a key_terms slide with 6 terms, each with a 20+ word definition")
    if quiz_slide:
        extras.append("a quiz slide: full-sentence question, 4 options, correct answer, 3-sentence explanation")
    if examples:
        example_note = "For each content slide include a vivid real-world example inside 'explanation'.\n"
    extras_text = ("Also include: " + "; ".join(extras) + ".\n") if extras else ""

    return f"""You are a {style.lower()} educator building a classroom PowerPoint for a teacher.

Topic : {topic}
Style : {style}   |   Difficulty : {difficulty}
Content slides (NOT counting title / objectives / summary): {num_content}

Curriculum material:
---
{context}
---

{extras_text}{example_note}
════════════════════════════════════════════════════════════════
STRICT CONTENT RULES — every slide must follow these exactly:
════════════════════════════════════════════════════════════════

1. bullets  → list of EXACTLY 5 items.
              Each bullet = a COMPLETE, INFORMATIVE sentence of 12-20 words.
              No vague phrases like "Key concept" or "Important point".
              Each must encode a concrete fact, definition, or cause-effect from the curriculum.

2. explanation → ONE paragraph of EXACTLY 60-80 words giving deeper context,
                 a worked example, or a real-world connection.
                 This appears as a panel on the slide — write it as the teacher's spoken elaboration.

3. speaker_notes → ONE paragraph of 100-130 words written in first person.
                   Include: an opening hook, 2-3 elaboration sentences, one
                   student-engagement question ("Can anyone...?"), and a transition.

4. title → max 8 words, title-case.

════════════════════════════════════════════════════════════════
JSON SCHEMA (return ONLY the raw array — no markdown, no prose)
════════════════════════════════════════════════════════════════

Content / objectives / summary slide:
{{
  "type": "content" | "objectives" | "summary",
  "title": "<8-word max>",
  "bullets": ["<12-20 word sentence>", "<12-20 word sentence>", ... x5],
  "explanation": "<60-80 word paragraph>",
  "speaker_notes": "<100-130 word script>"
}}

Title slide:
{{
  "type": "title",
  "title": "<topic name>",
  "subtitle": "<catchy 10-word descriptor>",
  "overview": "<3-sentence lesson overview paragraph, ~50 words>"
}}

Key-terms slide:
{{
  "type": "key_terms",
  "title": "Key Terms & Definitions",
  "terms": [
    {{"term": "<word>", "definition": "<20+ word definition sentence>"}},
    ... x6
  ]
}}

Quiz slide:
{{
  "type": "quiz",
  "title": "Quick Check",
  "question": "<full question sentence>",
  "options": ["<A>", "<B>", "<C>", "<D>"],
  "answer": "A"|"B"|"C"|"D",
  "explanation": "<3-sentence explanation of why the answer is correct>",
  "speaker_notes": "<50-word facilitation note>"
}}

SLIDE ORDER:
1. title
2. objectives   (5 bullets on what students will learn)
3-{2+num_content}. content  (one distinct sub-topic each)
{"next. key_terms" if key_terms else ""}
{"next. quiz" if quiz_slide else ""}
last.  summary   (5 takeaway bullets + next-steps explanation)
"""


def generate_slide_content(
    topic: str,
    num_slides: int = 8,
    teaching_style: str = "Visual",
    difficulty: str = "medium",
    include_examples: bool = True,
    include_key_terms: bool = True,
    include_quiz_slide: bool = False,
    teacher_name: str = "Teacher",
) -> list[dict]:
    """RAG + LLM → rich slide content list with robust parsing."""
    if vectorstore_exists():
        vs   = load_vectorstore()
        docs = retrieve_context(topic, k=8, vectorstore=vs)
        ctx  = "\n\n".join(
            f"[{d.metadata.get('title','Section')}]\n{d.page_content}" for d in docs
        )
    else:
        ctx = f"No curriculum uploaded. Use your comprehensive knowledge about: {topic}. Generate detailed, accurate, curriculum-grade content."

    num_content = max(3, num_slides - 3)
    prompt = _build_prompt(topic, ctx, num_content, teaching_style, difficulty,
                           include_examples, include_key_terms, include_quiz_slide)

    # Try up to 2 attempts
    for attempt in range(2):
        try:
            raw = get_llm(temperature=0.5).invoke(prompt).strip()
            parsed = _parse_llm_json(raw)
            if parsed and len(parsed) >= 3:
                # Validate each slide has at minimum a 'type'
                valid = []
                for s in parsed:
                    if isinstance(s, dict) and s.get("type"):
                        valid.append(s)
                if valid:
                    return valid
            print(f"[WARN] Attempt {attempt+1}: parsed {len(parsed) if parsed else 0} slides, retrying...")
        except Exception as e:
            print(f"[WARN] Attempt {attempt+1} failed: {e}")

    print(f"[WARN] LLM failed for slides on '{topic}', using fallback")
    return _fallback_slides(topic, teacher_name, num_content)


def _parse_llm_json(raw: str) -> list | None:
    """Parse LLM output to JSON array with multiple strategies."""
    # Strip markdown fences
    cleaned = re.sub(r"^```(?:json)?", "", raw, flags=re.IGNORECASE).strip()
    cleaned = re.sub(r"```$", "", cleaned).strip()

    # Strategy 1: Direct parse
    try:
        result = json.loads(cleaned)
        if isinstance(result, list):
            return result
    except json.JSONDecodeError:
        pass

    # Strategy 2: Find JSON array
    m = re.search(r"\[.*\]", cleaned, re.DOTALL)
    if m:
        try:
            result = json.loads(m.group(0))
            if isinstance(result, list):
                return result
        except json.JSONDecodeError:
            pass

    # Strategy 3: Fix common JSON issues (trailing commas, single quotes)
    fixed = cleaned
    fixed = re.sub(r",\s*([}\]])", r"\1", fixed)  # trailing commas
    fixed = fixed.replace("'", '"')  # single to double quotes
    m2 = re.search(r"\[.*\]", fixed, re.DOTALL)
    if m2:
        try:
            result = json.loads(m2.group(0))
            if isinstance(result, list):
                return result
        except json.JSONDecodeError:
            pass

    return None


def _fallback_slides(topic, teacher_name, num_content=4):
    """Rich fallback slides when LLM parsing fails."""
    slides = [
        {
            "type": "title", "title": topic,
            "subtitle": f"A comprehensive classroom lesson on {topic}",
            "overview": (
                f"This lesson introduces {topic} from first principles. "
                f"We explore key definitions, underlying mechanisms, and real-world applications. "
                f"By the end, students can explain core concepts and apply them to new problems."
            ),
        },
        {
            "type": "objectives", "title": "What You Will Learn Today",
            "bullets": [
                f"Define {topic} using precise, curriculum-grounded scientific language and terminology.",
                "Identify the main components of the system and describe each component's role clearly.",
                "Explain the underlying principles that govern the behaviour and dynamics of the system.",
                "Apply these principles to analyse practical, real-world scenarios and solve problems.",
                "Evaluate different approaches, compare their trade-offs, and justify optimal solutions.",
            ],
            "explanation": (
                f"These five objectives form a learning ladder — each one builds on the previous. "
                f"We start with vocabulary, move to understanding, and finish with application. "
                f"Use them as a checklist: if you can do all five by end of class, you're ready for assessment."
            ),
            "speaker_notes": (
                f"Good morning, everyone! Today's topic is {topic}. Before we dive in, look at these "
                f"five objectives on screen — they are your roadmap for the next hour. "
                f"We'll go from defining terms all the way to evaluating real-world trade-offs. "
                f"Can anyone tell me what they already know about {topic}? "
                f"Great — let's build on that. By the end of class I want you to test yourself: "
                f"can you explain each objective to a classmate without looking at your notes?"
            ),
        },
        {
            "type": "content", "title": f"Foundations of {topic}",
            "bullets": [
                f"{topic} is formally defined as a systematic framework used to describe and predict observable phenomena.",
                "The concept emerged from centuries of observation, hypothesis testing, and iterative scientific refinement.",
                "Its core principles provide a universal language shared across multiple disciplines and industries.",
                "Understanding the foundational model allows students to tackle both theoretical and applied problems.",
                "Modern technology, infrastructure, and scientific research all critically depend on these foundational ideas.",
            ],
            "explanation": (
                f"Think of {topic} as the grammar of a language — once you learn the rules, "
                f"you can construct and interpret an infinite number of sentences. "
                f"In engineering, these foundations underpin everything from bridge design to software architecture. "
                f"A solid grasp here will make every subsequent topic in this course significantly easier to absorb."
            ),
            "speaker_notes": (
                f"Let's start at the very beginning. {topic} is not just an abstract idea — it has shaped "
                f"human technology and science for decades. I want you to think of it as a lens: once you "
                f"understand it, you'll start recognising it in everything around you. "
                f"Can anyone spot an example from their own experience? Brilliant. Now let's look at "
                f"how it's formally defined, because precision of language matters enormously in this field."
            ),
        },
        {
            "type": "content", "title": f"Key Principles and Laws",
            "bullets": [
                f"The first principle states that {topic} operates under a framework of conservation and equilibrium.",
                "Mathematical relationships are used to quantify interactions between different system variables.",
                "Boundary conditions and constraints define the scope within which the principles are valid.",
                "Experimental validation has confirmed these principles across thousands of independent studies worldwide.",
                "When any principle is violated, it signals either a measurement error or a new discovery.",
            ],
            "explanation": (
                f"Principles are not arbitrary rules — they are distilled from thousands of experiments. "
                f"Each one tells you what WILL happen under specific conditions. "
                f"As an engineer or scientist, your job is to identify which principle applies to the problem at hand, "
                f"set up the correct equations, and solve for the unknown. Practice this workflow and it becomes intuitive."
            ),
            "speaker_notes": (
                f"This is the heart of {topic}. These principles are non-negotiable — nature follows them. "
                f"I always tell students: memorise the principle, but more importantly, understand WHY it works. "
                f"If you understand the 'why', you can derive everything else from scratch. "
                f"Can anyone think of a situation where Principle 1 would apply? Excellent. "
                f"Let's work through it together on the next slide."
            ),
        },
        {
            "type": "content", "title": f"Real-World Applications",
            "bullets": [
                f"{topic} is applied directly in modern engineering, medicine, computing, and environmental science.",
                "Industries use these concepts daily to design systems, optimise processes, and predict failures.",
                "Consumer technology from smartphones to electric vehicles relies on these exact principles.",
                "Research institutions are pushing the boundaries by applying these ideas to quantum and nano-scale systems.",
                "Students who master this topic gain a competitive advantage in internships, research, and interviews.",
            ],
            "explanation": (
                f"Applications make abstract theory tangible. When you see {topic} used in a real product or system, "
                f"the formulas stop being symbols on paper and become tools you can wield. "
                f"Take 5 minutes after class to look up one application that interests you — deepening that connection "
                f"will dramatically improve your recall during exams."
            ),
            "speaker_notes": (
                f"Now let's see why this matters beyond the classroom. {topic} is everywhere — from the device "
                f"in your pocket to the power grid that lights this room. I want you to feel that connection. "
                f"Pick ONE application that genuinely excites you and dig into it. "
                f"When you study for exams, that application will be your 'anchor', making everything else easier. "
                f"Can anyone share an application they've already encountered? Let's hear it."
            ),
        },
    ]

    # Add more content slides if needed
    extra_titles = [
        (f"Common Misconceptions About {topic}", "misconceptions"),
        (f"Problem-Solving Strategies", "strategies"),
        (f"Historical Development", "history"),
        (f"Advanced Topics & Extensions", "advanced"),
    ]
    for i in range(num_content - 3):
        if i >= len(extra_titles):
            break
        title, _ = extra_titles[i]
        slides.append({
            "type": "content", "title": title,
            "bullets": [
                f"Students commonly confuse this concept with closely related but fundamentally different ideas.",
                f"A structured approach — identify, plan, execute, verify — produces correct results more reliably.",
                f"Historical context shows how understanding of {topic} evolved through scientific revolutions.",
                f"Extension problems challenge students to apply principles in novel, unfamiliar scenarios.",
                f"Self-assessment questions help identify gaps before formal examinations and graded assignments.",
            ],
            "explanation": (
                f"This section addresses areas where students typically struggle most. "
                f"By explicitly naming common errors and providing counter-examples, we short-circuit bad habits "
                f"before they take root. Pay special attention to the differences between superficially similar concepts."
            ),
            "speaker_notes": (
                f"I've seen students make these mistakes year after year. Let's break the cycle today. "
                f"Look at each misconception carefully — if you recognise yourself, that's actually a good sign! "
                f"It means you're close to the right answer, you just need a small correction. "
                f"Can anyone tell me which of these surprised them? Great discussion — let's move on."
            ),
        })

    # Summary
    slides.append({
        "type": "summary", "title": "Recap & Next Steps",
        "bullets": [
            f"We established a precise definition of {topic} and its place within the broader field of study.",
            "Core principles and underlying mechanisms were explored with worked examples for deeper clarity.",
            "Key terminology was introduced to build the academic vocabulary essential for future assessments.",
            "Real-world applications illustrated why mastering this topic has practical, immediate career value.",
            "Next class we extend these ideas — review today's notes and attempt the practice questions tonight.",
        ],
        "explanation": (
            f"Today's session laid the critical foundation. Without a solid grasp of these concepts, "
            f"later topics will feel disconnected. Tonight, revisit your notes, highlight any gaps, "
            f"and complete 3-5 practice questions. Bring any unresolved doubts to the next session."
        ),
        "speaker_notes": (
            f"Excellent work today! Let's take 90 seconds to lock in what we covered. "
            f"We defined {topic}, explored its core principles, saw real-world applications, "
            f"and built up your vocabulary. For homework: read the relevant chapter, "
            f"attempt the five end-of-chapter questions, and note anything you're unsure about. "
            f"I'll start next class by answering your questions. Any questions right now? "
            f"If not, great job — see you next time!"
        ),
    })

    return slides


# ═══════════════════════════════════════════════════════════════════════════════
# LOW-LEVEL HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _bg(slide, color: RGBColor):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _rect(slide, l, t, w, h, color: RGBColor):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def _label(slide, text, l, t, w, h, sz=11, bold=False,
           color=None, align=PP_ALIGN.LEFT):
    """Single-run label — for short strings that will NOT wrap."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color or CLR_WHITE
    return tb


def _paragraph_frame(slide, l, t, w, h,
                      items: list[str],
                      sz=12,
                      color=None,
                      space_before_pt=7,
                      space_after_pt=3,
                      line_spacing_pt=17,
                      prefix="",
                      bold_prefix=False,
                      align=PP_ALIGN.LEFT):
    """
    Single textbox holding all items as separate paragraphs.
    Each paragraph uses fixed line-spacing so nothing overflows into the next.
    """
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space_before_pt)
        p.space_after  = Pt(space_after_pt)
        p.line_spacing = Pt(line_spacing_pt)

        if prefix:
            rp = p.add_run()
            rp.text = prefix + " "
            rp.font.size  = Pt(sz)
            rp.font.bold  = bold_prefix
            rp.font.color.rgb = CLR_ACCENT

        r = p.add_run()
        r.text = str(item)
        r.font.size = Pt(sz)
        r.font.color.rgb = color or CLR_WHITE

    return tb


def _notes(slide, text: str):
    if not text:
        return
    ns = slide.notes_slide.notes_text_frame
    ns.clear()
    p = ns.paragraphs[0]
    r = p.add_run()
    r.text = text


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE RENDERERS
# ═══════════════════════════════════════════════════════════════════════════════

# ── Title slide ───────────────────────────────────────────────────────────────
def _title_slide(prs, data, teacher, topic):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, CLR_BG_DARK)

    # Purple header block (top 38% of slide)
    _rect(sl, 0, 0, W, 2.85, CLR_BG_CARD)
    _rect(sl, 0, 0, W, 0.14, CLR_ACCENT)

    # Brand
    _label(sl, "MTRX-TriAxis  •  AI Classroom Assistant",
           0.45, 0.22, 7, 0.40, sz=11, bold=True, color=CLR_ACCENT)
    _rect(sl, 0.45, 0.78, 12.4, 0.025, CLR_ACCENT)

    # Topic title
    tb = sl.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(12.5), Inches(1.35))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r  = p.add_run(); r.text = data.get("title", topic)
    r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = CLR_WHITE

    # Subtitle
    sub_tb = sl.shapes.add_textbox(Inches(0.4), Inches(2.35), Inches(12.5), Inches(0.55))
    sub_tf = sub_tb.text_frame; sub_tf.word_wrap = True
    sp = sub_tf.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run(); sr.text = data.get("subtitle", f"A lesson by {teacher}")
    sr.font.size = Pt(16); sr.font.color.rgb = CLR_LGREY

    # Overview panel
    overview = data.get("overview", "")
    if overview:
        _rect(sl, 0.5, 3.15, 12.3, 2.55, CLR_BG_PANEL)
        _rect(sl, 0.5, 3.15, 0.075, 2.55, CLR_ACCENT2)
        _label(sl, "📋  Lesson Overview", 0.7, 3.22, 6, 0.4,
               sz=12, bold=True, color=CLR_ACCENT2)
        _paragraph_frame(sl, 0.7, 3.7, 12.0, 1.9,
                         [overview], sz=13, color=CLR_SILVER,
                         space_before_pt=0, line_spacing_pt=20)

    # Footer
    _rect(sl, 0, 6.9, W, 0.28, CLR_BG_PANEL)
    _label(sl, f"👩‍🏫  {teacher}", 0.5, 6.93, 6, 0.24, sz=10, color=CLR_GREY)
    _label(sl, f"📅  {date.today().strftime('%B %d, %Y')}",
           0.5, 6.93, 12.8, 0.24, sz=10, color=CLR_GREY, align=PP_ALIGN.RIGHT)

    _notes(sl, data.get("speaker_notes", ""))


# ── Content / Objectives / Summary slide ──────────────────────────────────────
#
# Layout (widths are for W=13.33"):
#
#  [0.0 ─ 13.33]  Top bar           0.00 – 0.14
#  [0.0 ─ 13.33]  Title band        0.14 – 1.05
#  [0.0 ─  8.55]  Bullet col        1.10 – 6.55   (8.3" wide, 5.45" tall)
#  [8.70 – 13.03] Explanation panel 1.10 – 6.55   (4.1" wide)
#   [0.0 ─ 13.33] Footer            6.65 – 7.50
#
def _content_slide(prs, data, num, total):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, CLR_BG_DARK)

    stype = data.get("type", "content")
    acc = {
        "objectives": CLR_ACCENT,
        "content":    CLR_ACCENT,
        "summary":    CLR_GREEN,
        "key_terms":  CLR_BLUE,
    }.get(stype, CLR_ACCENT)
    icon = {
        "objectives": "🎯",
        "content":    "📖",
        "summary":    "✅",
        "key_terms":  "📚",
    }.get(stype, "📖")

    # ── Header ──
    _rect(sl, 0, 0, W, 0.14, acc)
    _rect(sl, 0, 0.14, W, 0.91, CLR_BG_CARD)

    # Slide counter (top-right of header)
    _label(sl, f"{num} / {total}",
           0, 0.18, W - 0.35, 0.32, sz=9, color=CLR_GREY, align=PP_ALIGN.RIGHT)

    # Title
    title_tb = sl.shapes.add_textbox(
        Inches(0.4), Inches(0.20), Inches(11.8), Inches(0.72))
    title_tf = title_tb.text_frame; title_tf.word_wrap = True
    tp = title_tf.paragraphs[0]
    tr = tp.add_run()
    tr.text = f"{icon}  {data.get('title','')}"
    tr.font.size = Pt(24); tr.font.bold = True; tr.font.color.rgb = CLR_ACCENT2

    # ── LEFT: bullet list ──────────────────────────────────────────────────
    bullets    = data.get("bullets", [])
    bullet_col = 8.2    # width of bullet column
    bullet_l   = 0.40
    body_top   = 1.12
    body_h     = 5.38   # from 1.12 to 6.50

    _paragraph_frame(
        sl,
        l=bullet_l, t=body_top, w=bullet_col, h=body_h,
        items=bullets,
        sz=13,
        color=CLR_WHITE,
        space_before_pt=9,
        space_after_pt=2,
        line_spacing_pt=19,
        prefix="▸",
        bold_prefix=True,
    )

    # ── RIGHT: explanation panel ───────────────────────────────────────────
    expl = data.get("explanation", "")
    if expl:
        panel_l = 8.75
        panel_w = 4.20
        panel_t = body_top
        panel_h = body_h

        _rect(sl, panel_l, panel_t, panel_w, panel_h, CLR_BG_PANEL)
        _rect(sl, panel_l, panel_t, 0.07, panel_h, acc)

        _label(sl, "Context & Insight",
               panel_l + 0.16, panel_t + 0.12, panel_w - 0.2, 0.38,
               sz=10, bold=True, color=acc)

        _paragraph_frame(
            sl,
            l=panel_l + 0.16, t=panel_t + 0.56,
            w=panel_w - 0.22, h=panel_h - 0.65,
            items=[expl],
            sz=11.5,
            color=CLR_SILVER,
            space_before_pt=0,
            line_spacing_pt=18,
        )

    # ── Footer ─────────────────────────────────────────────────────────────
    _rect(sl, 0, 6.60, W, 0.06, acc)

    _notes(sl, data.get("speaker_notes", ""))


# ── Key Terms slide ────────────────────────────────────────────────────────────
#
# Two-column grid: 3 terms left, 3 terms right.
# Each term = bold label + definition paragraph in its own panel.
#
def _key_terms_slide(prs, data, num, total):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, CLR_BG_DARK)
    acc = CLR_BLUE

    _rect(sl, 0, 0, W, 0.14, acc)
    _rect(sl, 0, 0.14, W, 0.91, CLR_BG_CARD)
    _label(sl, f"{num} / {total}",
           0, 0.18, W - 0.35, 0.32, sz=9, color=CLR_GREY, align=PP_ALIGN.RIGHT)
    _label(sl, "📚  Key Terms & Definitions",
           0.4, 0.20, 10, 0.72, sz=24, bold=True, color=acc)

    raw_terms = data.get("terms", [])
    # Normalise: accept both {"term":…,"definition":…} and "Term: definition" strings
    terms = []
    for t in raw_terms:
        if isinstance(t, dict):
            terms.append((t.get("term", ""), t.get("definition", "")))
        else:
            parts = str(t).split(":", 1)
            terms.append((parts[0].strip(), parts[1].strip() if len(parts) > 1 else ""))

    # 3-column grid (up to 6 terms: 2 rows × 3 cols)
    cols     = 3
    col_w    = 4.05
    col_gap  = 0.25
    row_h    = 2.65
    start_l  = 0.4
    start_t  = 1.15

    for i, (term, defn) in enumerate(terms[:6]):
        col = i % cols
        row = i // cols
        l   = start_l + col * (col_w + col_gap)
        t   = start_t + row * (row_h + 0.15)

        _rect(sl, l, t, col_w, row_h, CLR_BG_PANEL)
        _rect(sl, l, t, 0.07, row_h, acc)

        # Term label
        _label(sl, term, l + 0.14, t + 0.12, col_w - 0.2, 0.42,
               sz=13, bold=True, color=acc)

        # Definition paragraph
        _paragraph_frame(
            sl,
            l=l + 0.14, t=t + 0.58,
            w=col_w - 0.22, h=row_h - 0.68,
            items=[defn],
            sz=11,
            color=CLR_SILVER,
            space_before_pt=0,
            line_spacing_pt=17,
        )

    _rect(sl, 0, 6.60, W, 0.06, acc)
    _notes(sl, data.get("speaker_notes", ""))


# ── Quiz slide ─────────────────────────────────────────────────────────────────
def _quiz_slide(prs, data, num, total):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, CLR_BG_DARK)
    acc = CLR_ORANGE

    _rect(sl, 0, 0, W, 0.14, acc)
    _rect(sl, 0, 0.14, W, 0.91, CLR_BG_CARD)
    _label(sl, f"{num} / {total}",
           0, 0.18, W - 0.35, 0.32, sz=9, color=CLR_GREY, align=PP_ALIGN.RIGHT)
    _label(sl, "🧠  Quick Check",
           0.4, 0.20, 10, 0.72, sz=24, bold=True, color=acc)

    # Question
    q_tb = sl.shapes.add_textbox(Inches(0.4), Inches(1.14), Inches(12.5), Inches(1.0))
    q_tf = q_tb.text_frame; q_tf.word_wrap = True
    qp = q_tf.paragraphs[0]; qp.line_spacing = Pt(20); qp.space_before = Pt(0)
    qr = qp.add_run(); qr.text = data.get("question", "")
    qr.font.size = Pt(16); qr.font.bold = True; qr.font.color.rgb = CLR_WHITE

    # Options — 2 × 2 grid
    opts   = data.get("options", [])
    labels = ["A", "B", "C", "D"]
    answer = data.get("answer", "A").strip().upper()

    positions = [
        (0.4,  2.28, 6.3),   # A left-top
        (6.9,  2.28, 6.0),   # B right-top
        (0.4,  3.60, 6.3),   # C left-bottom
        (6.9,  3.60, 6.0),   # D right-bottom
    ]
    for i, opt in enumerate(opts[:4]):
        lbl  = labels[i]
        l, t, pw = positions[i]
        is_a = (lbl == answer)
        bg   = RGBColor(0x0D, 0x2D, 0x18) if is_a else CLR_BG_PANEL
        clr  = CLR_GREEN if is_a else CLR_SILVER

        _rect(sl, l, t, pw, 1.05, bg)
        _rect(sl, l, t, 0.06, 1.05, clr)

        item_tb = sl.shapes.add_textbox(
            Inches(l + 0.14), Inches(t + 0.10), Inches(pw - 0.2), Inches(0.88))
        item_tf = item_tb.text_frame; item_tf.word_wrap = True
        ip = item_tf.paragraphs[0]; ip.line_spacing = Pt(17); ip.space_before = Pt(0)
        ir = ip.add_run(); ir.text = f"{lbl}.  {opt}"
        ir.font.size = Pt(13); ir.font.color.rgb = clr; ir.font.bold = is_a

    # Answer + explanation panel
    expl = data.get("explanation", "")
    _rect(sl, 0.4, 4.85, 12.5, 1.55, CLR_BG_PANEL)
    _rect(sl, 0.4, 4.85, 0.07, 1.55, CLR_GREEN)
    _label(sl, f"✅  Correct Answer: {answer}   —   Explanation",
           0.6, 4.88, 12.1, 0.4, sz=11, bold=True, color=CLR_GREEN)
    _paragraph_frame(
        sl,
        l=0.6, t=5.32, w=12.1, h=1.0,
        items=[expl or "Refer to your curriculum notes for details."],
        sz=11.5, color=CLR_SILVER,
        space_before_pt=0, line_spacing_pt=18,
    )

    _rect(sl, 0, 6.60, W, 0.06, acc)
    _notes(sl, data.get("speaker_notes", ""))


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN ASSEMBLER
# ═══════════════════════════════════════════════════════════════════════════════

def build_pptx(slides: list[dict], topic: str, teacher_name: str = "Teacher") -> bytes:
    """
    Convert slide data list → branded .pptx bytes.

    Returns bytes suitable for st.download_button().
    """
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    total = len(slides)

    for idx, slide_data in enumerate(slides, start=1):
        stype = slide_data.get("type", "content")
        if stype == "title":
            _title_slide(prs, slide_data, teacher_name, topic)
        elif stype == "key_terms":
            _key_terms_slide(prs, slide_data, idx, total)
        elif stype == "quiz":
            _quiz_slide(prs, slide_data, idx, total)
        else:
            _content_slide(prs, slide_data, idx, total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
